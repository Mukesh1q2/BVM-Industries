"""
BVM Industries — Company Presentation Generator
Creates a professional .pptx file with 14 slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0B, 0x0F, 0x17)
NAVY_LIGHT = RGBColor(0x0F, 0x15, 0x20)
BLUE       = RGBColor(0x2F, 0x8F, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xA9, 0xB3, 0xC2)
DARK_GRAY  = RGBColor(0x3A, 0x42, 0x50)

LOGO_PATH = r"c:\Users\crypt\Downloads\Ai projects\websites\Inprogress\BVM\app\bvm-next\public\new_logo.png"
OUT_PATH  = r"c:\Users\crypt\Downloads\Ai projects\websites\Inprogress\BVM\app\bvm-next\public\BVM_Industries_Presentation.pptx"

# ── Helpers ───────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # truly blank layout


def slide():
    """Add a new blank slide and paint it navy."""
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    return sl


def box(sl, left, top, width, height, text="", size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, bg=None, italic=False, wrap=True):
    """Add a text box; returns the text frame."""
    txb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tf


def rect(sl, left, top, width, height, color=BLUE, alpha=None):
    """Add a filled rectangle."""
    shape = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def logo(sl, left=0.25, top=0.15, height=0.55):
    """Place the BVM logo if the file exists."""
    if os.path.exists(LOGO_PATH):
        sl.shapes.add_picture(LOGO_PATH, Inches(left), Inches(top),
                              height=Inches(height))


def divider(sl, top, color=BLUE):
    """Thin horizontal rule."""
    rect(sl, 0.5, top, 12.33, 0.025, color)


def bullet_block(sl, items, left, top, width, height, icon="▸", icon_color=BLUE, text_color=GRAY, size=13):
    """Render a list of bullet items."""
    for i, item in enumerate(items):
        y = top + i * 0.38
        box(sl, left,        y, 0.3,   0.35, icon,  size, color=icon_color)
        box(sl, left + 0.28, y, width, 0.35, item,  size, color=text_color)


def tag(sl, left, top, text, bg=BLUE):
    """Small pill tag."""
    r = sl.shapes.add_shape(9, Inches(left), Inches(top), Inches(1.4), Inches(0.3))
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 4.8, 13.33, 2.7, NAVY_LIGHT)
rect(sl, 0, 4.75, 13.33, 0.05, BLUE)

logo(sl, left=0.6, top=0.7, height=1.0)

box(sl, 0.6, 2.05, 12, 0.7,
    "Precision Engineering for Aseptic Packaging Solutions",
    22, bold=False, color=BLUE, align=PP_ALIGN.LEFT, italic=True)

box(sl, 0.6, 2.9, 12, 1.1,
    "BVM INDUSTRIES", 52, bold=True, color=WHITE)

divider(sl, 4.15)

box(sl, 0.6, 4.95, 5, 0.4, "📍  Plot 774/496/46, Baddi, H.P. 173205", 11, color=GRAY)
box(sl, 0.6, 5.4,  5, 0.4, "✉️  sales@mybvm.in  |  service@mybvm.in", 11, color=GRAY)
box(sl, 0.6, 5.85, 5, 0.4, "📞  +91 7018231499", 11, color=GRAY)
box(sl, 0.6, 6.3,  5, 0.4, "🌐  mybvm.in", 11, color=BLUE)

box(sl, 8.5, 5.0, 4.5, 1.8,
    "Techno-Commercial\nPresentation\n2025–26",
    14, bold=False, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Who We Are
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "WHO WE ARE", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 10, 0.8, "BVM Industries", 36, bold=True, color=WHITE)
divider(sl, 1.72)

body = ("BVM Industries is a leading Indian manufacturer of aseptic packaging machinery. "
        "Founded on a commitment to precision engineering, we design and build F.F.S, B.F.S, "
        "and Euro Cap Sealing systems that meet the most rigorous pharmaceutical standards globally.")
box(sl, 0.5, 1.85, 7.5, 1.2, body, 13, color=GRAY)

stats = [
    ("8+",       "Years of Excellence"),
    ("26–50",    "Team Members"),
    ("₹5–25 Cr", "Annual Turnover"),
    ("Pan India","Service Coverage"),
]
for i, (val, label) in enumerate(stats):
    x = 0.5 + i * 3.2
    rect(sl, x, 3.3, 2.9, 1.6, NAVY_LIGHT)
    box(sl, x + 0.15, 3.45, 2.6, 0.7, val,   26, bold=True, color=BLUE,  align=PP_ALIGN.CENTER)
    box(sl, x + 0.15, 4.05, 2.6, 0.5, label, 10, bold=False, color=GRAY, align=PP_ALIGN.CENTER)

box(sl, 0.5, 5.2, 12, 0.4,
    "Proprietor: Rahul Kumar Singh  |  GST: 02GNLPS7342F1ZS  |  Location: Baddi, Himachal Pradesh",
    11, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Problem We Solve
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "THE CHALLENGE", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 10, 0.8, "Why Pharma Manufacturers Choose BVM", 30, bold=True, color=WHITE)
divider(sl, 1.72)

challenges = [
    ("Contamination Risk",   "Manual or semi-automated aseptic lines introduce contamination events that fail cGMP audits."),
    ("Import Dependency",    "Imported BFS/FFS lines carry 3–5× price premiums and 18-month lead times."),
    ("Compliance Burden",    "WHO-GMP, US FDA, and EU GMP documentation requires deep validation expertise."),
    ("Downtime Cost",        "Legacy machinery without modern HMI/PLC cannot produce audit trails — risking batch rejections."),
]
for i, (title, desc) in enumerate(challenges):
    x = 0.5 + (i % 2) * 6.3
    y = 2.1 + (i // 2) * 2.1
    rect(sl, x, y, 5.9, 1.8, NAVY_LIGHT)
    rect(sl, x, y, 0.06, 1.8, BLUE)
    box(sl, x + 0.2, y + 0.12, 5.5, 0.45, title, 14, bold=True, color=WHITE)
    box(sl, x + 0.2, y + 0.58, 5.5, 1.0,  desc,  11, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Product Portfolio Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "OUR PRODUCTS", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 10, 0.8, "Comprehensive Aseptic Packaging Solutions", 30, bold=True, color=WHITE)
divider(sl, 1.72)

products = [
    ("F.F.S Machines",      "Form-Fill-Seal systems for sterile pharmaceutical liquids."),
    ("B.F.S Machines",      "Blow-Fill-Seal for SVP & LVP parenterals with Class 100 aseptic chamber."),
    ("Euro Cap Sealing",    "400-series high-speed servo cap sealers for oral & ophthalmic lines."),
    ("Precision Moulds",    "BFS, PET and injection moulds in SS316L with micron-level tolerances."),
]
for i, (title, desc) in enumerate(products):
    x = 0.5 + i * 3.2
    rect(sl, x, 2.1, 2.9, 3.8, NAVY_LIGHT)
    rect(sl, x, 2.1, 2.9, 0.06, BLUE)
    box(sl, x + 0.15, 2.25, 2.6, 0.55, f"0{i+1}", 28, bold=True, color=BLUE, align=PP_ALIGN.LEFT)
    box(sl, x + 0.15, 2.82, 2.6, 0.6,  title, 13, bold=True, color=WHITE)
    box(sl, x + 0.15, 3.45, 2.6, 1.2,  desc,  11, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — F.F.S Machines
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "PRODUCT DEEP-DIVE", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 9, 0.8, "Form Fill Seal (F.F.S) Machines", 30, bold=True, color=WHITE)
divider(sl, 1.72)

box(sl, 0.5, 1.85, 7.5, 1.3,
    "Fully indigenous aseptic filling systems for injectables, ophthalmic, and respiratory products. "
    "Designed for continuous high-output parenteral production with Class 100 sterile chamber integration.",
    13, color=GRAY)

specs = ["Sterile chamber: Class 100 / ISO 5",
         "Container: 1 ml – 500 ml",
         "Materials: SS316L product contact parts",
         "Automation: Siemens S7 PLC + HMI",
         "Compliance: cGMP, WHO-GMP, 21 CFR Part 11"]
bullet_block(sl, specs, 0.5, 3.3, 6.0, 2.5)

tags = ["Injectables (SVP/LVP)", "Ophthalmic", "Respiratory"]
for i, t in enumerate(tags):
    tag(sl, 7.5 + i * 0.0, 1.85 + i * 0.45, t)

FFS_IMG = r"C:\Users\crypt\Downloads\Ai projects\websites\Inprogress\BVM\app\images\FSS-machine.png"
rect(sl, 8.8, 1.85, 4.0, 5.0, NAVY_LIGHT)
if os.path.exists(FFS_IMG):
    sl.shapes.add_picture(FFS_IMG, Inches(8.9), Inches(2.2), width=Inches(3.8))
else:
    box(sl, 8.85, 1.9, 3.9, 4.9, "[ FFS Machine Image ]", 12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — B.F.S Machines
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "PRODUCT DEEP-DIVE", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 9, 0.8, "Blow Fill Seal (B.F.S) Machines", 30, bold=True, color=WHITE)
divider(sl, 1.72)

box(sl, 0.5, 1.85, 7.5, 1.3,
    "Advanced cGMP-compliant systems for continuous aseptic filling of small and large volume "
    "parenterals. Our 7.0-Ton BFS platforms feature mirror-finish SS316L tooling and "
    "HGH45HA linear carriage assemblies.",
    13, color=GRAY)

highlights = [("7.0 Ton", "Machine capacity"), ("Class 100", "Sterile chamber"), ("SS316L", "All product contact parts"), ("24/7", "After-sales support")]
for i, (val, lbl) in enumerate(highlights):
    x = 0.5 + i * 1.78
    rect(sl, x, 3.3, 1.6, 1.2, NAVY_LIGHT)
    box(sl, x, 3.35, 1.6, 0.6, val, 18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    box(sl, x, 3.85, 1.6, 0.4, lbl, 9,  color=GRAY, align=PP_ALIGN.CENTER)

specs = ["Mould capacity: 1–6 cavity", "Volume: 1 ml – 1,000 ml",
         "IQ/OQ/PQ documentation included", "Compatible with WHO-GMP & EU GMP"]
bullet_block(sl, specs, 0.5, 4.7, 6.5, 2.0)

BFS_IMG = r"C:\Users\crypt\Downloads\Ai projects\websites\Inprogress\BVM\app\images\Blow Fill Seal (BFS).png"
rect(sl, 8.8, 1.85, 4.0, 5.0, NAVY_LIGHT)
if os.path.exists(BFS_IMG):
    sl.shapes.add_picture(BFS_IMG, Inches(8.9), Inches(2.2), width=Inches(3.8))
else:
    box(sl, 8.85, 1.9, 3.9, 4.9, "[ BFS Machine Image ]", 12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Euro Cap Sealing
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "PRODUCT DEEP-DIVE", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 9, 0.8, "Euro Cap Sealing — 400 Series", 30, bold=True, color=WHITE)
divider(sl, 1.72)

box(sl, 0.5, 1.85, 7.5, 1.3,
    "High-speed, servo-controlled cap sealing solutions for oral liquids and ophthalmic products. "
    "The 400 Series sealers deliver consistent repeatability across a wide range of bottle sizes "
    "with minimal changeover time.",
    13, color=GRAY)

specs = ["Servo-driven for precise torque control",
         "Handles multiple bottle/cap formats",
         "High-speed output: up to 400 caps/min",
         "Easy cap change tooling",
         "GMP-compliant design, easy to clean"]
bullet_block(sl, specs, 0.5, 3.3, 6.5, 2.5)

CAP_IMG = r"C:\Users\crypt\Downloads\Ai projects\websites\Inprogress\BVM\app\images\Cap-sealing-machine.png"
rect(sl, 8.8, 1.85, 4.0, 5.0, NAVY_LIGHT)
if os.path.exists(CAP_IMG):
    sl.shapes.add_picture(CAP_IMG, Inches(8.9), Inches(2.8), width=Inches(3.8))
else:
    box(sl, 8.85, 1.9, 3.9, 4.9, "[ Cap Sealer Image ]", 12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Precision Moulds
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "PRODUCT DEEP-DIVE", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 9, 0.8, "Precision Moulds & Tooling", 30, bold=True, color=WHITE)
divider(sl, 1.72)

mould_types = [
    ("BFS Moulds",       "Custom cavity designs for blow-fill-seal machines. Available in SS316L for sterile contact."),
    ("PET Bottle Moulds","High-clarity, consistent wall-thickness moulds for pharmaceutical PET containers."),
    ("Injection Moulds", "Multi-cavity injection moulds for caps, closures, and custom components."),
]
for i, (title, desc) in enumerate(mould_types):
    y = 2.0 + i * 1.65
    rect(sl, 0.5, y, 8.5, 1.4, NAVY_LIGHT)
    rect(sl, 0.5, y, 0.06, 1.4, BLUE)
    box(sl, 0.75, y + 0.1,  8.0, 0.5, title, 14, bold=True, color=WHITE)
    box(sl, 0.75, y + 0.58, 8.0, 0.7, desc,  11, color=GRAY)

box(sl, 0.5, 7.0, 12, 0.35,
    "Material: Grade SS316L  |  Tolerances: ±0.005 mm  |  Lead Time: 4–8 weeks  |  Custom configurations available",
    10, color=DARK_GRAY)

MOULD_IMG = r"C:\Users\crypt\Downloads\Ai projects\websites\Inprogress\BVM\app\images\mould-LVP.png"
if os.path.exists(MOULD_IMG):
    sl.shapes.add_picture(MOULD_IMG, Inches(9.2), Inches(2.2), width=Inches(3.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Turnkey Plant Services
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "SERVICES", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 9, 0.9, "Turnkey Plant Installation", 30, bold=True, color=WHITE)
divider(sl, 1.82)

box(sl, 0.5, 1.95, 8.0, 0.65,
    "Complete turnkey solutions for Parenteral Manufacturing Plants — from blueprint to batch production.",
    13, color=GRAY)

services = [
    "Conceptual Design & Master Planning",
    "Cleanroom Engineering & HVAC Systems",
    "Water Systems: PW / WFI / Clean Steam",
    "Sterile Filling & Lyophilization Lines",
    "Isolator & RABS Technology",
    "21 CFR Part 11 Automation & Compliance",
    "IQ/OQ/PQ Validation Documentation",
    "Audit Preparation & Regulatory Support",
]
bullet_block(sl, services[:4], 0.5, 2.8, 5.5, 3.0)
bullet_block(sl, services[4:], 6.5, 2.8, 6.0, 3.0)

standards = ["WHO-GMP", "cGMP", "EU GMP", "US FDA", "Data Integrity"]
for i, s in enumerate(standards):
    rect(sl, 0.5 + i * 2.55, 6.3, 2.3, 0.75, NAVY_LIGHT)
    box(sl, 0.5 + i * 2.55, 6.35, 2.3, 0.6, s, 10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Refurbishment Services
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "SERVICES", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 9, 0.8, "Machine Refurbishment & Modernisation", 30, bold=True, color=WHITE)
divider(sl, 1.72)

box(sl, 0.5, 1.85, 7.5, 0.65,
    "Complete overhaul and modernisation of legacy FFS/BFS machines. "
    "Upgrade to latest HMI/PLC standards — with resale and buyback options available.",
    13, color=GRAY)

steps = [
    ("01", "Condition Assessment", "Full mechanical, electrical and pneumatic audit of existing machine."),
    ("02", "HMI/PLC Upgrade",      "Replace legacy controls with Siemens S7 1500 and modern touchscreen HMI."),
    ("03", "Mechanical Overhaul",  "Replace worn carriages, seals, and product contact components with SS316L parts."),
    ("04", "Validation Support",   "IQ/OQ/PQ protocols for re-validated machine return to production."),
]
for i, (num, title, desc) in enumerate(steps):
    x = 0.5 + (i % 2) * 6.3
    y = 2.8 + (i // 2) * 2.0
    rect(sl, x, y, 5.9, 1.7, NAVY_LIGHT)
    box(sl, x + 0.15, y + 0.1,  0.7, 0.6, num,   20, bold=True, color=BLUE)
    box(sl, x + 0.85, y + 0.12, 4.8, 0.45, title, 13, bold=True, color=WHITE)
    box(sl, x + 0.85, y + 0.58, 4.8, 0.9,  desc,  11, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Industries We Serve
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "MARKETS", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 10, 0.8, "Industries We Serve", 30, bold=True, color=WHITE)
divider(sl, 1.72)

industries = [
    ("🏥", "Pharmaceutical",       "Sterile parenteral and ophthalmic packaging. Zero contamination to 21 CFR Part 11."),
    ("🍶", "Food & Beverage",      "High-throughput F.F.S machines for liquid to granular items. Hygienic & waste-efficient."),
    ("💄", "Cosmetic & Personal",  "Precision filling for lotions, creams, and serums with aesthetic seal integrity."),
    ("⚗️", "Chemical & Agro",      "Durable systems for aggressive substances. Leak-proof sealing per safety regulations."),
]
for i, (icon, title, desc) in enumerate(industries):
    x = 0.5 + i * 3.2
    rect(sl, x, 2.1, 2.9, 4.5, NAVY_LIGHT)
    box(sl, x + 0.15, 2.25, 2.6, 0.65, icon,  28, align=PP_ALIGN.CENTER)
    box(sl, x + 0.15, 2.9,  2.6, 0.55, title, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sl, x + 0.15, 3.55, 2.6, 2.8,  desc,  10, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Quality & Compliance
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "QUALITY ASSURANCE", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 10, 0.8, "Built for Compliance. Engineered for Trust.", 30, bold=True, color=WHITE)
divider(sl, 1.72)

pillars = [
    ("🛡️", "cGMP-Aligned Design",   "Cleanroom-ready layouts, CIP/SIP-friendly forms, and full validation documentation support from day one."),
    ("⏱️", "Timely Delivery",        "In-house machining, milestone tracking, and committed project schedules — no outsourcing delays."),
    ("🎧", "After-Sales Support",    "Dedicated spare parts inventory, remote troubleshooting, and on-site field service teams."),
    ("⚙️", "Siemens Automation",    "S7 1500 PLC + HMI control systems with 21 CFR Part 11 compliant audit trails and e-records."),
    ("🔬", "SS316L Materials",       "All product contact surfaces in premium Stainless Steel 316L — mirror polished to Ra ≤ 0.4 µm."),
    ("📄", "Validation Packages",    "IQ/OQ/PQ documentation prepared in house. Factory Acceptance Testing (FAT) included."),
]
for i, (icon, title, desc) in enumerate(pillars):
    x = 0.5 + (i % 3) * 4.27
    y = 2.1 + (i // 3) * 2.4
    rect(sl, x, y, 3.9, 2.1, NAVY_LIGHT)
    box(sl, x + 0.15, y + 0.1,  3.6, 0.5, f"{icon}  {title}", 12, bold=True, color=WHITE)
    box(sl, x + 0.15, y + 0.62, 3.6, 1.3, desc, 10, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Our Clients
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
logo(sl)
rect(sl, 0, 0, 0.07, 7.5, BLUE)
box(sl, 0.5, 0.4, 10, 0.5, "OUR CLIENTS", 11, bold=True, color=BLUE)
box(sl, 0.5, 0.85, 10, 0.8, "Trusted by Leading Pharmaceutical Manufacturers", 30, bold=True, color=WHITE)
divider(sl, 1.72)

box(sl, 0.5, 1.85, 12, 0.55,
    "Our machines are deployed across pharmaceutical, food, and chemical production facilities across India and beyond.",
    13, color=GRAY)

import glob
client_images = glob.glob(r"c:\Users\crypt\Downloads\Ai projects\websites\Inprogress\BVM\app\bvm-next\public\clients\*.png")

# Client logo placeholders (3×4 grid)
idx = 0
for row in range(3):
    for col in range(4):
        x = 0.5 + col * 3.2
        y = 2.7 + row * 1.5
        rect(sl, x, y, 2.9, 1.2, NAVY_LIGHT)
        if idx < len(client_images):
            try:
                sl.shapes.add_picture(client_images[idx], Inches(x + 0.2), Inches(y + 0.1), width=Inches(2.5))
            except Exception:
                pass
            idx += 1
        else:
            box(sl, x, y, 2.9, 1.2, f"Client {row*4+col+1}", 10, color=DARK_GRAY, align=PP_ALIGN.CENTER)

box(sl, 0.5, 7.1, 12, 0.3,
    "Reference list available on request  |  NDA maintained for all client relationships",
    10, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Contact & Next Steps
# ══════════════════════════════════════════════════════════════════════════════
sl = slide()
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 0, 13.33, 0.06, BLUE)
rect(sl, 0, 7.44, 13.33, 0.06, BLUE)

logo(sl, left=0.6, top=0.3, height=0.75)

box(sl, 0.6, 1.4, 12.0, 0.55, "LET'S BUILD SOMETHING TOGETHER", 11, bold=True, color=BLUE)
box(sl, 0.6, 1.95, 12.0, 1.1,  "Get in Touch With Our Team", 40, bold=True, color=WHITE)
divider(sl, 3.2)

contacts = [
    ("📞  Sales",    "+91 7018231499",       ""),
    ("✉️  Sales",    "sales@mybvm.in",       ""),
    ("✉️  Service",  "service@mybvm.in",     ""),
    ("📍  Address",  "Plot 774/496/46, Village Gullarwala", "Post & Teh. Baddi, Distt. Solan, H.P. 173205"),
    ("🌐  Website",  "mybvm.in",             ""),
]
for i, (label, primary, secondary) in enumerate(contacts):
    y = 3.45 + i * 0.73
    box(sl, 0.6,  y, 2.2, 0.5, label,   11, bold=True, color=BLUE)
    box(sl, 2.85, y, 5.5, 0.5, primary, 13, color=WHITE)
    if secondary:
        box(sl, 2.85, y + 0.33, 5.5, 0.38, secondary, 11, color=GRAY)

rect(sl, 9.0, 3.3, 3.8, 3.5, NAVY_LIGHT)
box(sl, 9.1, 3.45, 3.6, 0.5, "Next Steps", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(sl, 4.05)
next_steps = ["📋  Request a Techno-Commercial Offer", "🏭  Schedule a Factory Visit", "📐  Review G.A. Drawings", "🤝  Discuss Project Timeline"]
for i, step in enumerate(next_steps):
    box(sl, 9.1, 4.15 + i * 0.63, 3.6, 0.55, step, 10, color=GRAY)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"\n✅  Saved: {OUT_PATH}\n    Slides: {len(prs.slides)}\n")
